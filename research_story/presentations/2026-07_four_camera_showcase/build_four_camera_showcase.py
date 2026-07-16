#!/usr/bin/env python3
"""Build the four-camera warehouse reliability showcase deck.

The deck intentionally draws its map statistics from the checked-in day-zero
artifact rather than inventing illustrative coverage numbers.  It is a pitch
for the enlarged four-camera system and its reliability-fusion architecture;
the final scope slide distinguishes implemented components from empirical
claims that still require a four-camera campaign.
"""

from __future__ import annotations

import json
import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/four_camera_showcase_matplotlib")

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.colors import BoundaryNorm, ListedColormap
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
REPO = HERE.parents[2]  # UnembodiedNavigation/ (this file lives in research_story/presentations/<pkg>/)
MAP = REPO / "docs" / "assets" / "warehouse_full_4cam_map.png"
ARTIFACT = REPO / "paper_artifacts" / "gp" / "warehouse_full_4cam_dayzero_v1" / "camera_a_planner_with_four_camera_maps.npz"
MANIFEST = REPO / "paper_artifacts" / "gp" / "warehouse_full_4cam_dayzero_v1" / "prior_manifest.json"
# rendered assets + deck are regenerable bulk media -> logs/ (gitignored), never this tracked dir
ASSETS = REPO / "logs" / "studies" / "multicamera_commissioning_bigwarehouse" / "four_camera_showcase"
LIVE_VIEWS = ASSETS / "live_gazebo_views"
OUT = ASSETS / "Four_Camera_Warehouse_Reliability_Showcase.pptx"
ACTUAL_FIGURES = HERE / "full_story_walkthrough" / "07_real_commissioning_execution" / "figures"

W, H = 13.333, 7.5
FONT = "Aptos"
INK = RGBColor(0x12, 0x1B, 0x2A)
MUTED = RGBColor(0x5D, 0x68, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF9, 0xFC)
LINE = RGBColor(0xD8, 0xDE, 0xE8)
RED = RGBColor(0xC8, 0x19, 0x19)
BLUE = RGBColor(0x2F, 0x80, 0xED)
GREEN = RGBColor(0x21, 0x9A, 0x5B)
PURPLE = RGBColor(0x8D, 0x53, 0xC7)
ORANGE = RGBColor(0xED, 0x8A, 0x25)
TEAL = RGBColor(0x08, 0x8C, 0x98)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
PALE_BLUE = RGBColor(0xE9, 0xF3, 0xFF)
PALE_GREEN = RGBColor(0xE9, 0xF7, 0xEE)
PALE_ORANGE = RGBColor(0xFF, 0xF3, 0xE6)
PALE_PURPLE = RGBColor(0xF2, 0xEC, 0xFB)


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def image_aspect(path: Path) -> float:
    with Image.open(path) as image:
        return image.width / image.height


def add_text(
    slide,
    text: str,
    box: tuple[float, float, float, float],
    *,
    size: float = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    x, y, w, h = box
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_box(
    slide,
    box: tuple[float, float, float, float],
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LINE,
    radius: bool = True,
):
    x, y, w, h = box
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.7)
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_image(slide, path: Path, box: tuple[float, float, float, float], *, mode: str = "contain", border: RGBColor | None = None):
    x, y, w, h = box
    ratio = image_aspect(path)
    target = w / h
    if mode == "cover":
        dw, dh = (h * ratio, h) if ratio >= target else (w, w / ratio)
    else:
        dw, dh = (w, w / ratio) if ratio >= target else (h * ratio, h)
    picture = slide.shapes.add_picture(str(path), Inches(x + (w - dw) / 2), Inches(y + (h - dh) / 2), Inches(dw), Inches(dh))
    if border:
        outline = add_box(slide, box, fill=WHITE, line=border, radius=False)
        outline.fill.background()
    return picture


def add_title(slide, title: str, subtitle: str | None = None, section: str = "FOUR-CAMERA SHOWCASE") -> None:
    add_text(slide, title, (0.55, 0.30, 10.7, 0.48), size=25, bold=True)
    bar = add_box(slide, (0.57, 0.91, 1.02, 0.055), fill=RED, line=RED, radius=False)
    bar.line.fill.background()
    if subtitle:
        add_text(slide, subtitle, (0.57, 1.00, 10.75, 0.28), size=10.5, color=MUTED)
    add_text(slide, section, (10.4, 0.37, 2.35, 0.20), size=8.1, color=RED, bold=True, align=PP_ALIGN.RIGHT)


def add_footer(slide, number: int) -> None:
    add_text(slide, str(number), (0.55, 7.14, 0.35, 0.18), size=8.5, color=MUTED)
    add_text(slide, "Reliability-aware external-camera navigation", (0.88, 7.14, 4.0, 0.18), size=8.5, color=MUTED)
    add_text(slide, "warehouse_full_4cam", (10.2, 7.14, 2.55, 0.18), size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes.strip()


def add_metric(slide, label: str, value: str, detail: str, box: tuple[float, float, float, float], accent: RGBColor) -> None:
    x, y, w, h = box
    add_box(slide, box, fill=WHITE)
    stripe = add_box(slide, (x, y, 0.07, h), fill=accent, line=accent, radius=False)
    stripe.line.fill.background()
    add_text(slide, value, (x + 0.22, y + 0.18, w - 0.35, 0.42), size=23, color=accent, bold=True)
    add_text(slide, label, (x + 0.22, y + 0.67, w - 0.35, 0.23), size=10.6, bold=True)
    add_text(slide, detail, (x + 0.22, y + 0.98, w - 0.35, h - 1.08), size=8.8, color=MUTED)


def add_bullet_list(slide, items: list[str], box: tuple[float, float, float, float], *, size: float = 13.2, color: RGBColor = INK) -> None:
    x, y, w, h = box
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(7)
        paragraph.text = ""
        bullet = paragraph.add_run()
        bullet.text = "•  "
        bullet.font.name = FONT
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = RED
        run = paragraph.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = MUTED, width: float = 1.25) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True


def add_stage(slide, number: str, title: str, detail: str, box: tuple[float, float, float, float], accent: RGBColor) -> None:
    x, y, w, h = box
    add_box(slide, box, fill=WHITE)
    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.16), Inches(y + 0.18), Inches(0.38), Inches(0.38))
    disc.fill.solid()
    disc.fill.fore_color.rgb = accent
    disc.line.fill.background()
    add_text(slide, number, (x + 0.16, y + 0.265, 0.38, 0.16), size=8.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, (x + 0.64, y + 0.15, w - 0.78, 0.24), size=11, bold=True)
    add_text(slide, detail, (x + 0.64, y + 0.48, w - 0.78, h - 0.58), size=9.1, color=MUTED)


def draw_map_base(ax, xs: np.ndarray, ys: np.ndarray) -> None:
    ax.set_xlim(float(xs[0]), float(xs[-1]))
    ax.set_ylim(float(ys[0]), float(ys[-1]))
    ax.set_aspect("equal", adjustable="box")
    ax.set_xlabel("x [m]", fontsize=8)
    ax.set_ylabel("y [m]", fontsize=8)
    ax.tick_params(labelsize=7)
    ax.grid(color="#dce3eb", linewidth=0.6, alpha=0.8)
    for x, y, label, color in [(-6, -9.0, "A", "#2f80ed"), (-6, 9.0, "B", "#27ae60"), (6, -9.0, "C", "#9b51e0"), (6, 9.0, "D", "#f2994a")]:
        ax.scatter([x], [y], s=28, color=color, edgecolors="#17212f", linewidths=0.55, zorder=6)
        ax.text(x, y + (0.48 if y < 0 else -0.68), label, ha="center", va="center", fontsize=7, weight="bold", color="#17212f", zorder=7)


def render_reliability_atlas(data: np.lib.npyio.NpzFile) -> Path:
    xs, ys = data["xs"], data["ys"]
    entries = [
        ("Camera A", data["P_camera_A_map"], "magma"),
        ("Camera B", data["P_camera_B_map"], "magma"),
        ("Camera C", data["P_camera_C_map"], "magma"),
        ("Camera D", data["P_camera_D_map"], "magma"),
        ("4-camera union  p = 1 − ∏(1−pᵢ)", data["P_union_4cam_map"], "viridis"),
        ("Overlap count", data["coverage_count"], "overlap"),
    ]
    fig, axes = plt.subplots(2, 3, figsize=(13.2, 7.4), dpi=170, constrained_layout=True)
    for axis, (title, values, cmap) in zip(axes.ravel(), entries):
        if cmap == "overlap":
            palette = ListedColormap(["#f4f5f7", "#a8d5ff", "#4d9de0", "#1b6ca8", "#0b3c69"])
            norm = BoundaryNorm([-0.5, 0.5, 1.5, 2.5, 3.5, 4.5], palette.N)
            image = axis.imshow(values, origin="lower", extent=[xs[0], xs[-1], ys[0], ys[-1]], cmap=palette, norm=norm, interpolation="nearest")
            cbar = fig.colorbar(image, ax=axis, shrink=0.78, ticks=[0, 1, 2, 3, 4])
            cbar.ax.tick_params(labelsize=7)
            cbar.set_label("geometrically valid cameras", fontsize=7)
        else:
            image = axis.imshow(values, origin="lower", extent=[xs[0], xs[-1], ys[0], ys[-1]], cmap=cmap, vmin=0.0, vmax=1.0, interpolation="bilinear")
            cbar = fig.colorbar(image, ax=axis, shrink=0.78)
            cbar.ax.tick_params(labelsize=7)
            cbar.set_label("day-zero reliability", fontsize=7)
        draw_map_base(axis, xs, ys)
        axis.set_title(title, fontsize=10, weight="bold", pad=6)
    output = ASSETS / "dayzero_reliability_atlas.png"
    fig.savefig(output, facecolor="white")
    plt.close(fig)
    return output


def render_best_camera_map(data: np.lib.npyio.NpzFile) -> Path:
    xs, ys = data["xs"], data["ys"]
    ids = data["best_camera_id"].astype(str)
    codes = np.zeros(ids.shape, dtype=int)
    for index, name in enumerate(["camera_A", "camera_B", "camera_C", "camera_D"], start=1):
        codes[ids == name] = index
    cmap = ListedColormap(["#eef1f4", "#2f80ed", "#27ae60", "#9b51e0", "#f2994a"])
    fig, axes = plt.subplots(1, 2, figsize=(12.6, 5.4), dpi=170, constrained_layout=True)
    first = axes[0].imshow(codes, origin="lower", extent=[xs[0], xs[-1], ys[0], ys[-1]], cmap=cmap, vmin=0, vmax=4, interpolation="nearest")
    draw_map_base(axes[0], xs, ys)
    axes[0].set_title("Best day-zero camera by location", fontsize=12, weight="bold")
    cbar = fig.colorbar(first, ax=axes[0], ticks=[0, 1, 2, 3, 4], shrink=0.84)
    cbar.ax.set_yticklabels(["none", "A", "B", "C", "D"])
    cbar.ax.tick_params(labelsize=8)
    second = axes[1].imshow(data["P_best_4cam_map"], origin="lower", extent=[xs[0], xs[-1], ys[0], ys[-1]], cmap="viridis", vmin=0, vmax=1, interpolation="bilinear")
    draw_map_base(axes[1], xs, ys)
    axes[1].set_title("Best available reliability", fontsize=12, weight="bold")
    cbar = fig.colorbar(second, ax=axes[1], shrink=0.84)
    cbar.set_label("reliability", fontsize=8)
    cbar.ax.tick_params(labelsize=8)
    output = ASSETS / "best_camera_and_reliability.png"
    fig.savefig(output, facecolor="white")
    plt.close(fig)
    return output


def render_overlap_upgrade(data: np.lib.npyio.NpzFile, stats: dict) -> Path:
    xs, ys = data["xs"], data["ys"]
    coverage = data["coverage_count"]
    fig, axis = plt.subplots(figsize=(9.0, 5.2), dpi=170, constrained_layout=True)
    colors = ListedColormap(["#f5f6f8", "#c8e7ff", "#6daee6", "#2469a8", "#0d3c69"])
    image = axis.imshow(coverage, origin="lower", extent=[xs[0], xs[-1], ys[0], ys[-1]], cmap=colors, vmin=0, vmax=4, interpolation="nearest")
    draw_map_base(axis, xs, ys)
    axis.contour(xs, ys, coverage >= 2, levels=[0.5], colors=["#c81919"], linewidths=1.8)
    axis.text(0, 0.1, "red contour: ≥2-camera overlap", ha="center", va="center", fontsize=10, weight="bold", color="#8c1010", bbox={"boxstyle": "round,pad=0.25", "facecolor": "white", "edgecolor": "#c81919", "alpha": 0.88})
    axis.set_title("Overlap deliberately widened around the central handover corridor", fontsize=11.5, weight="bold", pad=20)
    cbar = fig.colorbar(image, ax=axis, shrink=0.82, ticks=[0, 1, 2, 3, 4])
    cbar.set_label("valid camera views", fontsize=9)
    cbar.ax.tick_params(labelsize=8)
    axis.text(0.02, 1.035, f"Union coverage: {stats['union']:.1%}    •    Multi-camera overlap: {stats['overlap']:.1%}", transform=axis.transAxes, ha="left", va="bottom", fontsize=8.8, weight="bold", color="#182230", bbox={"boxstyle": "round,pad=0.28", "facecolor": "white", "edgecolor": "#cfd8e3", "alpha": 0.95})
    output = ASSETS / "overlap_handover_corridor.png"
    fig.savefig(output, facecolor="white")
    plt.close(fig)
    return output


def _latest_live_frame(camera_id: str) -> Path:
    frames = sorted((LIVE_VIEWS / camera_id).glob("frame_*.png"))
    if not frames:
        raise FileNotFoundError(f"No Gazebo frame was captured for {camera_id!r}")
    return frames[-1]


def _font(size: int, *, bold: bool = False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def render_live_camera_montage() -> tuple[Path, Path]:
    selected = {
        "Camera A · south / west": _latest_live_frame("camera_A"),
        "Camera B · north / west": _latest_live_frame("camera_B"),
        "Camera C · south / east": _latest_live_frame("camera_C"),
        "Camera D · north / east": _latest_live_frame("camera_D"),
    }
    colors = ["#2f80ed", "#27ae60", "#9b51e0", "#f2994a"]
    canvas = Image.new("RGB", (1920, 1160), "#101b2a")
    draw = ImageDraw.Draw(canvas)
    draw.text((58, 38), "Live Gazebo views from the four wall-mounted cameras", fill="white", font=_font(42, bold=True))
    draw.text((60, 93), "Updated layout: camera columns at x = −6.0 m and x = +6.0 m", fill="#b7cae0", font=_font(22))
    panel_w, panel_h = 855, 481
    for index, ((label, frame), color) in enumerate(zip(selected.items(), colors)):
        col, row = index % 2, index // 2
        x, y = 60 + col * 930, 150 + row * 500
        image = Image.open(frame).convert("RGB")
        image.thumbnail((panel_w, panel_h), Image.Resampling.LANCZOS)
        px = x + (panel_w - image.width) // 2
        py = y + (panel_h - image.height) // 2
        draw.rectangle((x - 3, y - 3, x + panel_w + 3, y + panel_h + 3), fill=color)
        draw.rectangle((x, y, x + panel_w, y + panel_h), fill="#08111d")
        canvas.paste(image, (px, py))
        draw.rounded_rectangle((x + 18, y + 18, x + 365, y + 68), radius=10, fill="#111827")
        draw.text((x + 34, y + 29), label, fill="white", font=_font(21, bold=True))
    draw.text((60, 1124), "These are live RGB frames from Gazebo, not layout illustrations.", fill="#b7cae0", font=_font(20))
    montage = ASSETS / "live_four_camera_montage.png"
    canvas.save(montage, quality=95)
    overview = _latest_live_frame("overview")
    return montage, overview


def prepare_assets() -> tuple[dict, dict[str, Path]]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    if not MAP.exists() or not ARTIFACT.exists() or not MANIFEST.exists():
        raise FileNotFoundError("Build the four-camera world and day-zero prior before building the presentation.")
    manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
    summary = manifest["summary"]
    stats = {
        "union": float(summary["union_coverage_fraction"]),
        "overlap": float(summary["multi_covered_fraction"]),
        "mean_union": float(summary["mean_union_probability"]),
        "grid": f"{manifest['grid']['nx']} × {manifest['grid']['ny']}",
    }
    with np.load(ARTIFACT) as data:
        generated = {
            "map": MAP,
            "atlas": render_reliability_atlas(data),
            "best": render_best_camera_map(data),
            "overlap": render_overlap_upgrade(data, stats),
        }
    generated["live_montage"], generated["live_overview"] = render_live_camera_montage()
    actual = {
        "actual_routes": ACTUAL_FIGURES / "01_real_routes_and_observations.png",
        "actual_gp": ACTUAL_FIGURES / "02_actual_per_camera_gp_updates.png",
        "actual_gp_wide": ACTUAL_FIGURES / "02b_actual_gp_updates_wide.png",
        "actual_progress": ACTUAL_FIGURES / "03_actual_gp_learning_progress.png",
        "actual_overlap": ACTUAL_FIGURES / "04_actual_overlap_gate_C_D.png",
        "actual_execution": ACTUAL_FIGURES / "05_actual_algorithm_execution.png",
    }
    missing = [str(path) for path in actual.values() if not path.is_file()]
    if missing:
        raise RuntimeError(
            "Real commissioning figures are missing; run "
            "experiments/multicamera_commissioning_bigwarehouse/tools/"
            f"render_actual_commissioning_showcase.py first: {missing}"
        )
    generated.update(actual)
    return stats, generated


def make_presentation(stats: dict, images: dict[str, Path]) -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    blank = presentation.slide_layouts[6]

    # 1 — cover
    slide = presentation.slides.add_slide(blank)
    background = add_box(slide, (0, 0, W, H), fill=NAVY, line=NAVY, radius=False)
    background.line.fill.background()
    add_image(slide, images["map"], (6.35, 0.0, 6.98, 7.5), mode="cover")
    shade = add_box(slide, (5.9, 0, 1.2, 7.5), fill=NAVY, line=NAVY, radius=False)
    shade.fill.transparency = 20
    shade.line.fill.background()
    add_text(slide, "THE FOUR-CAMERA\nWAREHOUSE", (0.68, 1.06, 5.25, 1.40), size=31, color=WHITE, bold=True)
    add_box(slide, (0.72, 2.72, 1.18, 0.07), fill=RED, line=RED, radius=False).line.fill.background()
    add_text(slide, "A reliability-aware perception system\nthat sees around individual blind spots.", (0.72, 2.98, 4.9, 0.92), size=18, color=rgb("#D9E6F6"))
    add_text(slide, "24.5 × 20.5 m  ·  four calibrated views  ·  per-camera GP-ready priors  ·  conservative combination", (0.72, 4.48, 5.32, 0.35), size=9.4, color=rgb("#AFC4DC"), bold=True)
    add_text(slide, "Showcase presentation  |  July 2026", (0.72, 6.73, 3.8, 0.22), size=9.5, color=rgb("#AFC4DC"))
    add_notes(slide, "Open with the outcome: this is no longer a single camera bolted onto a warehouse. It is a four-view reliability system designed so that the robot can keep receiving credible position evidence when an individual view is weak or occluded.")

    # 2 — one picture
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "The upgrade, in one picture", "The warehouse is larger, the sensing is redundant, and reliability is now camera-specific.")
    add_image(slide, images["map"], (0.55, 1.46, 6.65, 5.30), mode="contain", border=LINE)
    add_metric(slide, "warehouse footprint", "24.5 × 20.5 m", "Dense rack blocks, props, a 4.5 m central aisle, and deliberately asymmetric shelf occlusions.", (7.55, 1.50, 2.48, 1.52), BLUE)
    add_metric(slide, "independent viewpoints", "4 cameras", "Two inward-looking south-wall views and two inward-looking north-wall views.", (10.25, 1.50, 2.48, 1.52), GREEN)
    add_metric(slide, "geometric union", f"{stats['union']:.1%}", "Of the planning grid is in at least one calibrated camera field of view on day zero.", (7.55, 3.34, 2.48, 1.52), PURPLE)
    add_metric(slide, "handover overlap", f"{stats['overlap']:.1%}", "Of the grid is seen by two or more cameras: the space to compare, confirm, and switch.", (10.25, 3.34, 2.48, 1.52), ORANGE)
    add_box(slide, (7.55, 5.20, 5.18, 1.15), fill=PALE_BLUE, line=rgb("#B8D5F5"))
    add_text(slide, "Core idea", (7.78, 5.40, 1.0, 0.22), size=10.5, color=BLUE, bold=True)
    add_text(slide, "Do not ask “is the camera available?” Ask “which camera is credible here, right now?”", (7.78, 5.67, 4.62, 0.45), size=14, bold=True)
    add_footer(slide, 2)
    add_notes(slide, "Use the layout to make the scale tangible. The four-camera system is not redundancy for its own sake: the tall, asymmetric rack segments make different camera views fail in different places. The overlap region gives the system somewhere to validate a source switch rather than blindly jumping between cameras.")

    # 3 — overlap design
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "Overlap is designed in—not left to chance", "The two camera columns are shifted inward from ±7.5 m to ±6.0 m to expand the useful handover band.")
    add_image(slide, images["overlap"], (0.58, 1.38, 7.55, 5.72), mode="contain", border=LINE)
    add_metric(slide, "previous layout", "37.9%", "of the grid had multi-camera geometric coverage.", (8.55, 1.70, 3.75, 1.24), MUTED)
    add_metric(slide, "updated layout", f"{stats['overlap']:.1%}", "multi-camera overlap after moving the columns inward.", (8.55, 3.14, 3.75, 1.24), RED)
    add_box(slide, (8.55, 4.84, 3.75, 1.42), fill=PALE_ORANGE, line=rgb("#F4D2AB"))
    add_text(slide, "+4.3 percentage points", (8.80, 5.08, 3.2, 0.27), size=17, color=ORANGE, bold=True)
    add_text(slide, "More overlap means more places to compare observations, recover from a drop-out, or pass localization to a better view.", (8.80, 5.48, 3.05, 0.54), size=10.3, color=INK)
    add_footer(slide, 3)
    add_notes(slide, "This was a deliberate layout refinement. Moving the camera columns inward increases the overlap from 37.9% to 42.2%, while the union coverage rises to 99.2%. The red contour is the region where at least two calibrated views can contribute.")

    # 4 — live camera montage
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "Fresh Gazebo views from all four live camera topics", "Each panel is a current RGB frame from the updated warehouse—not a schematic or a synthetic crop.")
    add_image(slide, images["live_montage"], (0.52, 1.35, 12.30, 5.67), mode="contain", border=LINE)
    add_footer(slide, 4)
    add_notes(slide, "Show the audience that the cameras are real Gazebo sensors looking through the actual shelf geometry. The views are intentionally different: each sees a different combination of rack faces, aisles, props, shadows, and occlusions. That is exactly why they need separate reliability models.")

    # 5 — full-facility overview
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "A high Gazebo view reveals the whole operating field", "This is a presentation-only overview sensor: it shows the facility, but it is not a fifth localization camera.")
    add_image(slide, images["live_overview"], (0.55, 1.48, 8.25, 5.28), mode="contain", border=LINE)
    add_box(slide, (9.18, 1.70, 3.45, 1.18), fill=PALE_BLUE, line=rgb("#B8D5F5"))
    add_text(slide, "Whole-facility context", (9.45, 1.96, 2.92, 0.24), size=14, color=BLUE, bold=True)
    add_text(slide, "The rack geometry, central aisle, no-go boundaries, staging areas, and robot are visible in one live image.", (9.45, 2.29, 2.85, 0.33), size=9.7, color=MUTED)
    add_box(slide, (9.18, 3.18, 3.45, 1.18), fill=PALE_GREEN, line=rgb("#B9DFC7"))
    add_text(slide, "Visual-only sensor", (9.45, 3.44, 2.92, 0.24), size=14, color=GREEN, bold=True)
    add_text(slide, "Mounted above the centre at (0, 0, 26 m). Its topic is excluded from the four-camera GP, fusion, and planner paths.", (9.45, 3.77, 2.85, 0.37), size=9.7, color=MUTED)
    add_box(slide, (9.18, 4.66, 3.45, 1.18), fill=PALE_ORANGE, line=rgb("#F4D2AB"))
    add_text(slide, "Why it helps the pitch", (9.45, 4.92, 2.92, 0.24), size=14, color=ORANGE, bold=True)
    add_text(slide, "It lets the audience immediately understand the scale and spatial context behind the four oblique camera views.", (9.45, 5.25, 2.85, 0.37), size=9.7, color=MUTED)
    add_footer(slide, 5)
    add_notes(slide, "This view is deliberately labelled presentation-only. It is not used by the four-camera system, so it never inflates the coverage or reliability claims. Its role is to make the world legible in a talk: you can see the complete warehouse before we return to the hard, oblique sensing views.")

    # 6 — day-zero atlas
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "Every camera starts with its own day-zero reliability map", "Six maps from the calibrated artifact: four camera-specific fields, their union, and the overlap count.")
    add_image(slide, images["atlas"], (0.46, 1.32, 12.42, 5.68), mode="contain", border=LINE)
    add_footer(slide, 6)
    add_notes(slide, "This is the heart of the upgrade. We do not use one generic camera reliability map. Each camera has a separate spatial field that knows its own field of view, scale, distance, image-border margin, and ground incidence. The union map makes clear that reliability is a spatial system property, not a property of an individual frame.")

    # 7 — initial GP semantics
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "An initial GP is a starting hypothesis, not false certainty", "The day-zero fields use calibration only; experience then corrects the parts that geometry cannot see.")
    add_box(slide, (0.65, 1.52, 3.68, 4.82), fill=PALE_BLUE, line=rgb("#B8D5F5"))
    add_text(slide, "What the initial field knows", (0.95, 1.85, 3.05, 0.26), size=14, color=BLUE, bold=True)
    add_bullet_list(slide, ["Camera pose and intrinsics", "Field-of-view membership", "Distance, obliquity, and pixel scale", "Image-border margin and target height"], (0.95, 2.35, 3.0, 2.32), size=12.4)
    add_text(slide, "No training data. No ground-truth labels.", (0.95, 5.40, 2.95, 0.30), size=11.8, color=BLUE, bold=True)
    add_box(slide, (4.83, 1.52, 3.68, 4.82), fill=PALE_ORANGE, line=rgb("#F4D2AB"))
    add_text(slide, "What it cannot know", (5.13, 1.85, 3.05, 0.26), size=14, color=ORANGE, bold=True)
    add_bullet_list(slide, ["Shelf occlusion and shadows", "Detector failure modes", "Lighting changes or calibration drift", "A camera that is stale or degraded"], (5.13, 2.35, 3.0, 2.32), size=12.4)
    add_text(slide, "Geometry can over-predict visibility—never prove it.", (5.13, 5.32, 2.95, 0.42), size=11.4, color=ORANGE, bold=True)
    add_box(slide, (9.00, 1.52, 3.68, 4.82), fill=PALE_GREEN, line=rgb("#B9DFC7"))
    add_text(slide, "How the GP becomes safe", (9.30, 1.85, 3.05, 0.26), size=14, color=GREEN, bold=True)
    add_bullet_list(slide, ["Low-strength prior: it can be overturned", "Camera-specific observations update only that map", "Predict mean and uncertainty between samples", "Plan conservatively when uncertainty remains high"], (9.30, 2.35, 3.0, 2.32), size=12.4)
    add_text(slide, "Learned evidence refines the prior; it does not contaminate the other cameras.", (9.30, 5.16, 3.00, 0.53), size=11.2, color=GREEN, bold=True)
    add_footer(slide, 7)
    add_notes(slide, "Be precise on the term initial GP. It is a calibrated reliability prior over the floor, not a detector-trained claim of visibility. The entire point of a low-strength prior is that just a few misses can overturn an optimistic geometric prediction behind a rack.")

    # 8 — actual uncertainty-stamped collection
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "Actual four-camera commissioning runs", "Handover and A/C-overlap passes recorded source identity, detections, misses, timestamps, and noisy operational odometry.")
    add_box(slide, (0.62, 1.55, 2.75, 3.95), fill=PALE_BLUE, line=rgb("#B8D5F5"))
    add_text(slide, "1  Execute route families", (0.92, 1.90, 2.15, 0.25), size=14, color=BLUE, bold=True)
    add_bullet_list(slide, ["A–D single-source regions", "Adjacent-camera overlap corridors", "Both directions, speeds, and lateral offsets", "One camera-gap crossing"], (0.92, 2.42, 2.08, 1.92), size=10.8)
    add_text(slide, "Routes are intentionally not invented: the current config remains empty until each is validated in this world.", (0.92, 4.62, 2.08, 0.56), size=9.1, color=BLUE, bold=True)
    add_arrow(slide, 3.55, 3.44, 4.12, 3.44, color=MUTED, width=1.6)
    add_box(slide, (4.30, 1.55, 4.18, 3.95), fill=PALE_PURPLE, line=rgb("#D8C8EE"))
    add_text(slide, "2  Preserve the operational record", (4.62, 1.90, 3.55, 0.25), size=14, color=PURPLE, bold=True)
    fields = [("camera", "A / B / C / D"), ("observation", "pixel pose + detector quality"), ("time", "timestamp + frame age"), ("state", "ŝₜ and Σₛ,ₜ (pose covariance)")]
    for index, (label, value) in enumerate(fields):
        y = 2.43 + index * 0.53
        add_box(slide, (4.64, y, 3.45, 0.38), fill=WHITE, line=rgb("#D8C8EE"), radius=False)
        add_text(slide, label, (4.78, y + 0.10, 0.92, 0.14), size=8.1, color=PURPLE, bold=True)
        add_text(slide, value, (5.83, y + 0.09, 2.06, 0.16), size=8.7, color=INK)
    add_text(slide, "A miss is evidence too. The source and uncertainty remain attached to every record.", (4.64, 4.73, 3.36, 0.34), size=9.6, color=PURPLE, bold=True)
    add_arrow(slide, 8.67, 3.44, 9.24, 3.44, color=MUTED, width=1.6)
    add_box(slide, (9.42, 1.55, 3.28, 3.95), fill=PALE_GREEN, line=rgb("#B9DFC7"))
    add_text(slide, "3  Keep evaluation separate", (9.72, 1.90, 2.68, 0.25), size=14, color=GREEN, bold=True)
    add_bullet_list(slide, ["Operational stream drives the GP and manager", "Ground truth is held for scoring only", "Held-out routes prevent a route replay from looking like generalisation", "Per-camera CSVs share a timeline without sharing labels"], (9.72, 2.42, 2.52, 1.92), size=10.4)
    add_text(slide, "This is the firewall that makes a later GP and fusion claim credible.", (9.72, 4.68, 2.52, 0.38), size=9.3, color=GREEN, bold=True)
    add_box(slide, (1.48, 5.92, 10.38, 0.56), fill=PALE_ORANGE, line=rgb("#F4D2AB"))
    add_text(slide, "PROTOCOL / DATA PENDING  ·  The next generated visual is an executed-route map with covariance ellipses and per-camera coverage—not a decorative route sketch.", (1.72, 6.10, 9.90, 0.18), size=10.2, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, (0.46, 1.36, 12.40, 5.47), fill=WHITE, line=LINE, radius=False)
    add_image(slide, images["actual_routes"], (0.48, 1.38, 12.36, 5.43), mode="contain")
    add_footer(slide, 8)
    add_notes(slide, "This is an executed result, not a route sketch. The long south-to-north traverse and dedicated A/C overlap pass recorded all four detector streams together with an independent noisy-odometry stream. The figure keeps detections and misses separate. The 0.10 m state covariance floor is explicit because the current encoder-noise publisher has zero covariance entries.")

    # 9 — fit individual GPs from the real records
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "Four actual GP updates—one per camera", "Each expected-kernel posterior uses only its own operational records and calibrated day-zero prior.")
    stages = [
        ("m₀,c(s)", "Calibration-only prior", "Pose, FOV, range, pixel scale, and image-border margin.", BLUE),
        ("D0 / D1", "Uncertain-input records", "Detections and misses, each stamped with ŝ and Σₛ.", PURPLE),
        ("μc(s), σc(s)", "Camera-specific posterior", "Fit camera A, B, C, and D independently with frozen choices.", GREEN),
        ("hold-out", "Earn the map claim", "NLL, MAE, calibration, and false-high-trust exposure.", ORANGE),
    ]
    for index, (symbol, title, detail, color) in enumerate(stages):
        x = 0.66 + index * 3.17
        add_box(slide, (x, 1.77, 2.78, 2.30), fill=WHITE, line=color)
        add_text(slide, symbol, (x + 0.20, 2.09, 2.35, 0.32), size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, (x + 0.20, 2.64, 2.35, 0.26), size=11.8, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, (x + 0.25, 3.18, 2.28, 0.46), size=8.7, color=MUTED, align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            add_arrow(slide, x + 2.80, 2.90, x + 3.08, 2.90, color=MUTED, width=1.45)
    add_box(slide, (0.98, 4.65, 11.38, 1.20), fill=rgb("#F7F9FC"), line=LINE)
    add_text(slide, "The learned-GP figure set", (1.28, 4.91, 2.20, 0.23), size=12.2, color=NAVY, bold=True)
    add_text(slide, "For A–D: observations + pose covariance  |  posterior mean  |  posterior standard deviation  |  held-out calibration card", (3.46, 4.89, 8.45, 0.31), size=12.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Unvisited cells remain visibly uncertain; no camera may borrow another camera's learned failures.", (1.52, 5.42, 10.20, 0.19), size=10.1, color=MUTED, align=PP_ALIGN.CENTER)
    add_box(slide, (3.08, 6.23, 7.18, 0.43), fill=PALE_ORANGE, line=rgb("#F4D2AB"))
    add_text(slide, "DATA PENDING  ·  The current atlas is day-zero prior only; these posterior plots appear only after D0/D1 passes.", (3.30, 6.36, 6.75, 0.14), size=8.9, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, (0.40, 1.26, 12.54, 5.69), fill=WHITE, line=LINE, radius=False)
    add_image(slide, images["actual_gp_wide"], (0.42, 1.28, 12.50, 5.65), mode="contain")
    add_footer(slide, 9)
    add_notes(slide, "These are actual pilot fits: 60–62 aligned records per camera across two executed routes. The figure makes the right boundary visible: observed paths change the calibrated prior, while unvisited space retains high posterior uncertainty. This establishes implementation and the direction of learned change, not route-disjoint generalisation.")

    # 10 — actual overlap gate
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "C/D overlap gate: real pilot pass", "Synchronized Camera C/D observations are checked for agreement before any combination claim.")
    overlap_steps = [
        ("1", "Find geometric overlap", "Current day-zero map: 42.2% of the grid has ≥2 valid views.", BLUE),
        ("2", "Collect D2 pairs", "Record only contemporaneous projected observations in adjacent-camera corridors.", PURPLE),
        ("3", "Pass the agreement gate", "≥30 held-out synchronized pairs per claimed edge; ≤10% spatial outliers.", ORANGE),
        ("4", "Choose the safe action", "Select first; combine only when consistency and covariance checks pass.", GREEN),
    ]
    for index, (number, title, detail, color) in enumerate(overlap_steps):
        x = 0.62 + index * 3.18
        add_stage(slide, number, title, detail, (x, 1.72, 2.80, 1.46), color)
        if index < len(overlap_steps) - 1:
            add_arrow(slide, x + 2.82, 2.43, x + 3.08, 2.43, color=MUTED, width=1.45)
    outputs = [("Overlap graph", "which physical camera pairs were tested"), ("Sync timeline", "timestamp offset and pair count"), ("Disagreement map", "where compatible views diverge")]
    for index, (title, detail) in enumerate(outputs):
        x = 1.22 + index * 3.83
        add_box(slide, (x, 4.10, 3.20, 1.22), fill=WHITE, line=LINE)
        add_text(slide, title, (x + 0.20, 4.37, 2.78, 0.22), size=13.0, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, (x + 0.25, 4.76, 2.70, 0.26), size=9.3, color=MUTED, align=PP_ALIGN.CENTER)
    add_box(slide, (1.37, 5.93, 10.58, 0.55), fill=PALE_ORANGE, line=rgb("#F4D2AB"))
    add_text(slide, "DATA PENDING  ·  Until this gate passes, the system may show availability overlap but must not claim validated measurement fusion.", (1.66, 6.11, 10.00, 0.18), size=10.0, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, (0.46, 1.40, 12.38, 5.38), fill=WHITE, line=LINE, radius=False)
    add_image(slide, images["actual_overlap"], (0.48, 1.42, 12.34, 5.34), mode="contain")
    add_footer(slide, 10)
    add_notes(slide, "This is genuine pilot overlap evidence. Three C/D pairs were synchronized within the strict 50 ms tolerance. Their mean disagreement was 0.247 m and each stayed below the 0.30 m gate. The correct conclusion is deliberately narrow: the pilot C/D edge passed, while a campaign-level fusion claim still requires the preregistered 30 held-out pairs.")

    # 11 — distinct GPs and availability union
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "Four camera-specific reliability maps stay separate", "The day-zero availability union visualizes coverage; live measurement combination needs the D2 consistency gate.")
    labels = [("GP A", BLUE, "camera A detections / misses"), ("GP B", GREEN, "camera B detections / misses"), ("GP C", PURPLE, "camera C detections / misses"), ("GP D", ORANGE, "camera D detections / misses")]
    for i, (name, color, detail) in enumerate(labels):
        x = 0.72 + i * 3.12
        add_box(slide, (x, 1.68, 2.68, 1.18), fill=WHITE, line=color)
        add_text(slide, name, (x + 0.20, 1.91, 2.25, 0.25), size=16, color=color, bold=True)
        add_text(slide, detail, (x + 0.20, 2.29, 2.22, 0.28), size=9.4, color=MUTED)
        add_arrow(slide, x + 1.34, 2.92, x + 1.34, 3.42, color=color)
    add_box(slide, (1.28, 3.60, 10.75, 1.20), fill=rgb("#F7F9FC"), line=LINE)
    add_text(slide, "At robot position s:", (1.60, 3.86, 2.0, 0.22), size=12.4, bold=True)
    add_text(slide, "p₁(s), p₂(s), p₃(s), p₄(s)  →  select the most credible source or combine consistent sources", (3.10, 3.82, 8.30, 0.32), size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "p_available(s) = 1 − ∏ᵢ [1 − pᵢ(s)]", (3.10, 5.25, 7.20, 0.42), size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Availability union improves coverage; measurement fusion still requires synchronized observations to agree.", (2.25, 5.88, 8.86, 0.30), size=12.2, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)
    add_notes(slide, "The formula is an availability union, not blind coordinate averaging. At runtime, the system can select the best source or sequentially fuse only consistent map observations after the D2 gate. We retain source identity all the way through, so Camera A’s occlusion tells us nothing automatically about Camera B.")

    # 12 — best-camera map
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "The day-zero map turns four viewpoints into a spatial sensing policy", "At each location, choose the most credible current camera; use overlap to corroborate or combine only after commissioning.")
    add_image(slide, images["best"], (0.54, 1.43, 7.74, 5.28), mode="contain", border=LINE)
    add_box(slide, (8.70, 1.68, 3.75, 1.22), fill=PALE_BLUE, line=rgb("#B8D5F5"))
    add_text(slide, "1. Select", (8.98, 1.94, 1.10, 0.23), size=14, color=BLUE, bold=True)
    add_text(slide, "Pick the best available source by reliability, detector quality, and freshness.", (8.98, 2.28, 3.05, 0.29), size=9.8, color=MUTED)
    add_box(slide, (8.70, 3.16, 3.75, 1.22), fill=PALE_GREEN, line=rgb("#B9DFC7"))
    add_text(slide, "2. Confirm", (8.98, 3.42, 1.40, 0.23), size=14, color=GREEN, bold=True)
    add_text(slide, "Where two views overlap, check that map estimates are spatially compatible.", (8.98, 3.76, 3.05, 0.29), size=9.8, color=MUTED)
    add_box(slide, (8.70, 4.64, 3.75, 1.22), fill=PALE_PURPLE, line=rgb("#D8C8EE"))
    add_text(slide, "3. Fuse or defer", (8.98, 4.90, 1.90, 0.23), size=14, color=PURPLE, bold=True)
    add_text(slide, "Fuse consistent updates; otherwise preserve the belief and increase uncertainty.", (8.98, 5.24, 3.05, 0.29), size=9.8, color=MUTED)
    add_footer(slide, 12)
    add_notes(slide, "The left map is an intuitive spatial policy. It says which camera has the best day-zero reliability at every floor location. The actual live decision incorporates observation quality and freshness as well; we do not execute a static zone switch blindly.")

    # 13 — actual algorithm execution
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "Actual end-to-end algorithm execution", "Weakly trusted pilot regions are withheld rather than injected as unverified corrections.")
    stages = [
        ("A–D", "Four independent RGB streams", "Dedicated camera topics and detector instances.", BLUE),
        ("GP", "Per-camera reliability", "Independent day-zero priors and online updates.", GREEN),
        ("✓", "Select / sequentially fuse", "Use reliability, detector score, age, and overlap agreement.", PURPLE),
        ("R", "Handover-adjusted covariance", "Keep mean; inflate uncertainty when the switch is unproven.", ORANGE),
        ("B", "Belief + planner", "The robot receives a calibrated external-camera correction.", RED),
    ]
    for index, (symbol, title, detail, color) in enumerate(stages):
        x = 0.55 + index * 2.55
        add_box(slide, (x, 2.05, 2.12, 2.22), fill=WHITE, line=color)
        disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.78), Inches(2.31), Inches(0.56), Inches(0.56))
        disc.fill.solid(); disc.fill.fore_color.rgb = color; disc.line.fill.background()
        add_text(slide, symbol, (x + 0.78, 2.45, 0.56, 0.20), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, (x + 0.18, 3.05, 1.76, 0.37), size=11.2, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, (x + 0.18, 3.52, 1.76, 0.47), size=8.9, color=MUTED, align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            add_arrow(slide, x + 2.15, 3.16, x + 2.48, 3.16, color=MUTED, width=1.5)
    add_box(slide, (1.28, 5.20, 10.77, 0.78), fill=rgb("#F7F9FC"), line=LINE)
    add_text(slide, "Key property: the planning interface stays the same—only the credibility of the camera measurement becomes spatially and source aware.", (1.58, 5.45, 10.15, 0.23), size=12.2, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, (0.42, 1.46, 12.49, 5.09), fill=WHITE, line=LINE, radius=False)
    add_image(slide, images["actual_execution"], (0.44, 1.48, 12.45, 5.05), mode="contain")
    add_footer(slide, 13)
    add_notes(slide, "This is the actual replay chain. Two routes delivered 242 aligned records, which fitted four expected-kernel posteriors plus a pooled diagnostic GP. The C/D pilot overlap passed. In the hysteretic replay, all 57 candidate frames remained below the configured 0.45 spatial-trust release threshold, so zero corrections were released. That safe-defer outcome is an integration and safety result, not a claim of closed-loop improvement.")

    # 14 — handover
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "A handover is a measurement event, not a camera toggle", "Switching source is safe only when the incoming update is supported by the overlap evidence.")
    add_stage(slide, "1", "Detect a source change", "The chosen camera is different from the last accepted source.", (0.72, 1.72, 2.30, 1.33), BLUE)
    add_stage(slide, "2", "Ask whether overlap agrees", "Compare A/B/C/D map observations where the two views should see the same robot.", (3.20, 1.72, 2.30, 1.33), GREEN)
    add_stage(slide, "3", "Check observation health", "Account for detector quality, timestamp age, missing frames, and reliability loss.", (5.68, 1.72, 2.30, 1.33), PURPLE)
    add_stage(slide, "4", "Adjust covariance", "Accept a confirmed handover—or inflate R until the new source proves itself.", (8.16, 1.72, 2.30, 1.33), ORANGE)
    add_stage(slide, "5", "Correct the belief", "The same filter receives a measurement whose trust matches the evidence.", (10.64, 1.72, 2.00, 1.33), RED)
    for x in [3.04, 5.52, 8.00, 10.48]:
        add_arrow(slide, x, 2.38, x + 0.12, 2.38)
    add_box(slide, (1.40, 3.74, 10.55, 1.55), fill=PALE_ORANGE, line=rgb("#F4D2AB"))
    add_text(slide, "If the views disagree, the robot should not average its way into confidence.", (1.78, 4.07, 9.80, 0.30), size=17, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Keep the position estimate honest, reduce the influence of the unsupported observation, and wait for a consistent update.", (1.78, 4.57, 9.80, 0.30), size=12.5, color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, 14)
    add_notes(slide, "This is the safety argument. Camera switching is a source change in a probabilistic measurement system. We assess agreement, staleness, quality, and whether overlap confirmation exists. The correct response to uncertainty is covariance inflation—not pretending that an abrupt source switch is equally trustworthy.")

    # 15 — walkthrough
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "What a four-camera traverse looks like", "A conceptual walkthrough of how the system uses redundancy across a shelf-occluded route.")
    walkthrough = [
        ("South approach", "A + C", "Two south-wall cameras establish and cross-check the initial external observation.", BLUE),
        ("Tall-shelf shadow", "C leads", "A can lose confidence locally; C is evaluated independently rather than inheriting A’s failure.", PURPLE),
        ("Central handover band", "compare / fuse", "Multiple calibrated views create a chance to confirm the next source before relying on it.", GREEN),
        ("North exit", "B + D", "The northern pair becomes primary, carrying a source-specific reliability history.", ORANGE),
    ]
    for index, (place, active, description, color) in enumerate(walkthrough):
        x = 0.72 + index * 3.10
        add_box(slide, (x, 1.78, 2.72, 3.98), fill=WHITE, line=color)
        add_text(slide, f"0{index + 1}", (x + 0.22, 2.02, 0.44, 0.25), size=13, color=color, bold=True)
        add_text(slide, place, (x + 0.22, 2.44, 2.20, 0.26), size=14, bold=True)
        tag = add_box(slide, (x + 0.22, 2.98, 1.66, 0.36), fill=color, line=color)
        tag.line.fill.background()
        add_text(slide, active, (x + 0.28, 3.08, 1.54, 0.14), size=8.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, description, (x + 0.22, 3.72, 2.22, 1.25), size=10.2, color=MUTED)
        if index < 3:
            add_arrow(slide, x + 2.74, 3.72, x + 3.00, 3.72, color=color, width=1.7)
    add_box(slide, (1.35, 6.18, 10.65, 0.54), fill=PALE_BLUE, line=rgb("#B8D5F5"))
    add_text(slide, "The key is continuity of credible localization—not a promise that every camera sees everywhere.", (1.60, 6.34, 10.15, 0.18), size=11.7, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 15)
    add_notes(slide, "Describe this as a system walkthrough, not an empirical result. The visual makes clear why camera-specific GPs matter: a local failure by A need not propagate to C, and the system obtains an evidence-based bridge into the northern cameras through the overlap band.")

    # 16 — run/demo
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "The showcase is reproducible from the world to the reliability artifact", "Everything on the previous slides is generated from the canonical world, its four camera calibrations, and a planner-compatible artifact.")
    commands = [
        ("1", "Generate the world", "python3 scripts/geometry_visibility/make_warehouse_full.py", BLUE),
        ("2", "Build four day-zero maps", "python3 scripts/geometry_visibility/build_full4cam_planner_prior.py", GREEN),
        ("3", "Launch four detector streams", "ros2 launch experiments warehouse_full4cam_commissioning.launch.py", PURPLE),
        ("4", "Collect, then replay", "design routes → export A–D records → evaluate the gated policies", ORANGE),
    ]
    for index, (number, label, code, color) in enumerate(commands):
        x = 0.72 + index * 3.10
        add_box(slide, (x, 1.82, 2.72, 2.64), fill=WHITE, line=color)
        add_text(slide, number, (x + 0.22, 2.12, 0.30, 0.25), size=14, color=color, bold=True)
        add_text(slide, label, (x + 0.22, 2.54, 2.15, 0.31), size=13.2, bold=True)
        add_box(slide, (x + 0.22, 3.16, 2.25, 0.76), fill=rgb("#F2F5F8"), line=rgb("#E1E6EC"), radius=False)
        add_text(slide, code, (x + 0.32, 3.30, 2.02, 0.48), size=7.5, color=NAVY)
    add_box(slide, (0.86, 5.16, 11.64, 1.05), fill=PALE_GREEN, line=rgb("#B9DFC7"))
    add_text(slide, "Presentation deliverables", (1.18, 5.42, 1.86, 0.24), size=11.3, color=GREEN, bold=True)
    add_text(slide, "large four-camera world  →  camera-specific priors  →  union / best / overlap maps  →  selection and handover interface", (3.07, 5.39, 8.95, 0.29), size=14.0, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 16)
    add_notes(slide, "This slide is practical proof that the presentation is tied to buildable assets rather than decorative illustrations. The artifact contains the four individual maps, union and best-camera maps, coverage count, and source IDs. The extension launch provides isolated detector streams on the four image topics.")

    # 17 — honest scope
    slide = presentation.slides.add_slide(blank)
    add_title(slide, "What this showcase establishes—and what the campaign must still prove", "Strong system engineering is explicit about the boundary between implemented capability and empirical performance.")
    add_box(slide, (0.76, 1.58, 5.72, 4.90), fill=PALE_GREEN, line=rgb("#B9DFC7"))
    add_text(slide, "Implemented in the four-camera showcase", (1.10, 1.94, 4.92, 0.26), size=15.0, color=GREEN, bold=True)
    add_bullet_list(slide, ["Canonical large warehouse with four isolated image streams", "Per-camera calibrated day-zero reliability maps", "Union, best-camera, and overlap fields in one artifact", "Source-specific reliability-provider interface", "Selection, sequential fusion, and handover-covariance logic"], (1.10, 2.48, 4.82, 2.90), size=12.0)
    add_box(slide, (6.86, 1.58, 5.72, 4.90), fill=PALE_ORANGE, line=rgb("#F4D2AB"))
    add_text(slide, "What a four-camera campaign must measure", (7.20, 1.94, 4.92, 0.26), size=15.0, color=ORANGE, bold=True)
    add_bullet_list(slide, ["Detector reliability and calibration per physical view", "Localization continuity through real shelf occlusions", "Whether fusion outperforms selection under drop-outs", "Handover error, latency, and covariance calibration", "Closed-loop route-level benefit versus a single camera"], (7.20, 2.48, 4.82, 2.90), size=12.0)
    add_text(slide, "The architecture is now ready to make those comparisons meaningful.", (1.38, 6.66, 10.60, 0.24), size=13.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 17)
    add_notes(slide, "End with credibility. The system layout, maps, and fusion interfaces are implemented. The system should not yet be described as having a measured four-camera localization improvement until a campaign produces the per-camera detector evidence and closed-loop handover result. This is a strength: we have built the right experiment, not skipped it.")

    # 18 — close
    slide = presentation.slides.add_slide(blank)
    backdrop = add_box(slide, (0, 0, W, H), fill=NAVY, line=NAVY, radius=False)
    backdrop.line.fill.background()
    add_text(slide, "FROM ONE CAMERA FIELD\nTO A RELIABLE SENSING NETWORK", (0.75, 1.12, 7.2, 1.20), size=30, color=WHITE, bold=True)
    add_box(slide, (0.78, 2.62, 1.24, 0.065), fill=RED, line=RED, radius=False).line.fill.background()
    add_text(slide, "Four cameras give coverage.\nPer-camera evidence earns calibrated GP expectations.\nConservative combination protects the robot through uncertainty.", (0.78, 2.95, 7.2, 1.35), size=18, color=rgb("#D8E6F7"))
    add_box(slide, (8.66, 1.35, 3.63, 3.95), fill=WHITE, line=WHITE)
    add_text(slide, f"{stats['union']:.1%}", (8.98, 1.90, 2.95, 0.55), size=31, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "geometric union coverage", (8.98, 2.54, 2.95, 0.25), size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"{stats['overlap']:.1%}", (8.98, 3.34, 2.95, 0.55), size=31, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "multi-camera handover overlap", (8.98, 3.98, 2.95, 0.25), size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Next: use the campaign to turn this sensing architecture into measured closed-loop evidence.", (0.78, 6.58, 10.8, 0.25), size=12.4, color=rgb("#B6CBE1"))
    add_notes(slide, "Close on the simple transformation: coverage is not enough, because each viewpoint has different failure modes. The upgrade adds four independent sources, four reliability models, an overlap-aware handover mechanism, and a clear empirical path to validate the closed-loop benefit.")

    return presentation


def main() -> None:
    stats, images = prepare_assets()
    presentation = make_presentation(stats, images)
    presentation.save(OUT)
    print(f"Wrote {OUT}")
    print(f"Generated reliability visuals in {ASSETS}")
    print(f"Union coverage: {stats['union']:.1%}; multi-camera overlap: {stats['overlap']:.1%}")


if __name__ == "__main__":
    main()
